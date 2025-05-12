import streamlit as st
import base64
import pptx
import json
from pptx.util import Inches, Pt
import os
from openai import OpenAI
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Load language files
def load_languages():
    try:
        with open('language.json', 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        st.error("Language file not found. Please ensure 'language.json' exists.")
        st.stop()
    except json.JSONDecodeError:
        st.error("Invalid JSON format in language file.")
        st.stop()

LANGUAGES = load_languages()

# Check for API Key
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
if not DEEPSEEK_API_KEY:
    st.error("API Key not found in .env file. Please check and try again.")
    st.stop()

# Setup DeepSeek API Client
client = OpenAI(
    api_key=DEEPSEEK_API_KEY,
    base_url="https://api.deepseek.com"
)

# Font settings
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def analyze_content(topic, content, lang):
    """Analyze content and return structured analysis for slide generation"""
    texts = LANGUAGES[lang]
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": texts.get("system_prompt_analysis", 
                "You are an expert presentation content analyzer. Analyze the provided content and return: "
                "1. Key themes/topics (for slide titles) "
                "2. Important points (for slide content) "
                "3. Recommended structure "
                "Format as JSON with keys: themes, points, recommendations")},
            {"role": "user", "content": texts.get("user_prompt_analysis", 
                "Analyze this content for presentation about {topic}:\n\n{content}").format(
                topic=topic, content=content)}
        ],
        response_format={ "type": "json_object" },
        stream=False
    )
    
    try:
        return json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        return {"themes": [], "points": [], "recommendations": ""}

def generate_slide_titles(topic, count, lang, analysis=None):
    """Generate slide titles incorporating analysis results"""
    texts = LANGUAGES[lang]
    
    base_prompt = texts["user_prompt_titles"].format(topic=topic, count=count)
    
    if analysis:
        prompt = (f"{base_prompt}\n\nAnalysis Results:\n"
                 f"Key Themes: {', '.join(analysis.get('themes', []))}\n"
                 f"Recommended Structure: {analysis.get('recommendations', '')}\n\n"
                 "Please generate slide titles that incorporate these themes.")
    else:
        prompt = base_prompt
    
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": texts["system_prompt_titles"]},
            {"role": "user", "content": prompt}
        ],
        stream=False
    )
    return [title.strip() for title in response.choices[0].message.content.split("\n") if title.strip()]

def generate_slide_content(slide_title, lang, analysis=None):
    """Generate slide content incorporating analysis results"""
    texts = LANGUAGES[lang]
    
    base_prompt = texts["user_prompt_content"].format(slide_title=slide_title)
    
    if analysis:
        relevant_points = "\n".join([
            f"- {point}" for point in analysis.get('points', []) 
            if any(keyword.lower() in slide_title.lower() 
                 for keyword in analysis.get('themes', []))
        ])
        
        prompt = (f"{base_prompt}\n\nRelevant Points from Analysis:\n"
                 f"{relevant_points}\n\n"
                 "Incorporate these points where appropriate.")
    else:
        prompt = base_prompt
    
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": texts["system_prompt_content"]},
            {"role": "user", "content": prompt}
        ],
        stream=False
    )
    return response.choices[0].message.content

def create_presentation(topic, slide_titles, slide_contents):
    """Create PowerPoint presentation with the generated content"""
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Content slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Font styling
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # Create directory if not exists
    os.makedirs("generated_ppt", exist_ok=True)
    prs.save(f"generated_ppt/{topic}_presentation.pptx")

def get_ppt_download_link(topic, lang):
    """Generate download link for the presentation"""
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">{LANGUAGES[lang]["download"]}</a>'

def main():
    # Initialize session state for analysis results
    if 'analysis_results' not in st.session_state:
        st.session_state.analysis_results = None
    
    # Language selection
    lang = st.selectbox("Select Language / 选择语言 / เลือกภาษา", list(LANGUAGES.keys()))
    texts = LANGUAGES[lang]

    st.title(texts["title"])
    st.subheader(texts["subtitle"])
    st.markdown('<style>h1{color: orange; text-align: center;}</style>', unsafe_allow_html=True)
    st.markdown('<style>h3{color: pink; text-align: center;}</style>', unsafe_allow_html=True)

    # Main input
    topic = st.text_input(texts["topic_prompt"])
    
    # Slide count selector
    slide_count = st.slider(texts["slide_count"], 1, 10, 5)
    
    # Content analysis section
    with st.expander(texts.get("content_analysis", "Content Analysis")):
        content_to_analyze = st.text_area(
            texts.get("content_analysis_prompt", "Enter content to analyze (optional):"),
            placeholder=texts.get("content_analysis_placeholder", "Paste relevant content that AI should consider"),
            height=150
        )
        
        if st.button(texts.get("analyze_button", "Analyze Content")) and content_to_analyze:
            with st.spinner(texts.get("analyzing", "Analyzing content...")):
                st.session_state.analysis_results = analyze_content(topic, content_to_analyze, lang)
                
                # Display analysis summary
                st.success(texts.get("analysis_complete", "Analysis complete!"))
                st.subheader(texts.get("key_themes", "Key Themes"))
                st.write(", ".join(st.session_state.analysis_results.get("themes", [])))
                
                st.subheader(texts.get("recommendations", "Recommended Structure"))
                st.write(st.session_state.analysis_results.get("recommendations", ""))
    
    # Custom slides option
    with st.expander(texts["custom_slides"]):
        custom_titles = st.text_area(
            texts["custom_slides"],
            placeholder=texts["custom_placeholder"],
            height=150
        )
    
    # Generate presentation button
    if st.button(texts["generate_btn"]):
        if not topic:
            st.warning(texts["no_topic_warning"])
        else:
            with st.spinner(texts["generating"]):
                try:
                    # Get analysis results if available
                    analysis = st.session_state.analysis_results if 'analysis_results' in st.session_state else None
                    
                    # Process custom titles if provided
                    if custom_titles.strip():
                        slide_titles = [title.strip() for title in custom_titles.split("\n") if title.strip()]
                        # Ensure we don't exceed requested slide count
                        slide_titles = slide_titles[:slide_count]
                        # If user provided fewer titles than requested count, generate the rest
                        if len(slide_titles) < slide_count:
                            remaining_count = slide_count - len(slide_titles)
                            generated_titles = generate_slide_titles(topic, remaining_count, lang, analysis)
                            slide_titles.extend(generated_titles)
                    else:
                        # Generate all titles automatically
                        slide_titles = generate_slide_titles(topic, slide_count, lang, analysis)
                    
                    # Generate content for each slide
                    slide_contents = []
                    for title in slide_titles:
                        content = generate_slide_content(title, lang, analysis)
                        slide_contents.append(content)
                    
                    # Create presentation
                    create_presentation(topic, slide_titles, slide_contents)
                    
                    st.success(texts["success"])
                    st.markdown(get_ppt_download_link(topic, lang), unsafe_allow_html=True)
                    
                except Exception as e:
                    st.error(f"{texts.get('error', 'Error')}: {str(e)}")

if __name__ == "__main__":
    main()
